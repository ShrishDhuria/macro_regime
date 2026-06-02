"""Phase 6 - generate the methodology PowerPoint deck."""
from __future__ import annotations
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from data.storage import load_panel

ROOT = Path(__file__).resolve().parent.parent
REPORTS = ROOT / "reports"

NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
RED   = RGBColor(0xC1, 0x12, 0x1F)
GRAY  = RGBColor(0x4A, 0x4A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_slide(prs, title_text, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tb.text_frame.text = title_text
    p = tb.text_frame.paragraphs[0]
    p.font.size, p.font.bold, p.font.color.rgb = Pt(28), True, NAVY
    if section:
        sb = slide.shapes.add_textbox(Inches(10), Inches(0.35), Inches(2.8), Inches(0.4))
        sb.text_frame.text = section
        sp = sb.text_frame.paragraphs[0]
        sp.font.size, sp.font.color.rgb, sp.alignment = Pt(11), GRAY, PP_ALIGN.RIGHT
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY; line.line.fill.background()
    return slide


def add_text(slide, text, left, top, width, height, size=14, bold=False, color=GRAY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.text = text
    for p in tb.text_frame.paragraphs:
        p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(size), bold, color, align
    return tb


def add_bullets(slide, items, left, top, width, height, size=14):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "-   " + item
        p.font.size, p.font.color.rgb = Pt(size), GRAY
        p.space_after = Pt(8)
    return tb


def add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return None
    if width and height:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def add_table(slide, df, left, top, width, height, size=10, header_size=11):
    rows, cols = df.shape[0] + 1, df.shape[1]
    ts = slide.shapes.add_table(rows, cols, left, top, width, height); t = ts.table
    for j, col in enumerate(df.columns):
        c = t.cell(0, j); c.text = str(col)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs:
            p.font.color.rgb, p.font.bold, p.font.size = WHITE, True, Pt(header_size)
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            v = df.iloc[i, j]
            txt = (f"{v:.4f}" if abs(v) < 10 else f"{v:.2f}") if isinstance(v, float) else str(v)
            c = t.cell(i + 1, j); c.text = txt
            for p in c.text_frame.paragraphs:
                p.font.size, p.font.color.rgb = Pt(size), GRAY
    return ts


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.6))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2.9), Inches(11.3), Inches(1.6))
    tb.text_frame.text = "Macro Regime &\nCross-Asset Risk Intelligence Platform"
    for p in tb.text_frame.paragraphs:
        p.font.size, p.font.bold, p.font.color.rgb = Pt(40), True, NAVY
    add_text(s, "An institutional-grade research framework for European markets",
              Inches(1), Inches(4.7), Inches(11.3), Inches(0.6), size=18, color=GRAY)
    add_text(s, "Shrish Dhuria  |  ESSEC Master in Finance  |  May 2026",
              Inches(1), Inches(6.7), Inches(11.3), Inches(0.4), size=14, color=GRAY)


def slide_summary(prs):
    s = add_slide(prs, "Project Summary", "Overview")
    add_text(s, "Mission", Inches(0.5), Inches(1.3), Inches(12), Inches(0.4), size=16, bold=True, color=NAVY)
    add_text(s, "Build a regime-aware macro-financial framework that detects regimes across European "
              "cross-asset markets, computes regime-conditional risk metrics, and tests tactical "
              "allocation strategies under realistic walk-forward backtest conditions.",
              Inches(0.5), Inches(1.7), Inches(12), Inches(1.2), size=14)
    add_text(s, "Six layers", Inches(0.5), Inches(3.0), Inches(12), Inches(0.4), size=16, bold=True, color=NAVY)
    add_bullets(s, [
        "Data spine - 14 series across equity, FX, commodity, rates, macro (2005-2026, weekly)",
        "Feature library - 63 features across 7 categories with explicit lookahead-bias enforcement",
        "Regime detection - 3-state Gaussian Hidden Markov Model on 5 macro features",
        "Risk engine - VaR, ES, Cornish-Fisher, regime-conditional metrics + Excel parallel build",
        "Tactical allocation - 6 strategies, walk-forward backtest, realistic frictions",
        "Stress testing - empirical-beta-based scenario analysis"
    ], Inches(0.5), Inches(3.4), Inches(12), Inches(3.5), size=14)


def slide_data(prs):
    s = add_slide(prs, "Data Infrastructure", "Phase 1")
    add_bullets(s, [
        "14 macro and market series via yfinance, FRED, ECB Data Portal",
        "Aligned to weekly Friday-close grid; ~1,114 weekly observations",
        "Source diversification with FRED-mirror fallbacks (institutional reliability pattern)",
        "Automated freshness check - flags series stale >90 days (caught OECD CLI drift day one)",
        "Lookahead-bias policy: macro releases lagged by publication delay before model use"
    ], Inches(0.5), Inches(1.3), Inches(12), Inches(2.5), size=14)
    cov = pd.DataFrame({"Class": ["Equity", "FX", "Commodity", "Volatility", "Rates", "Macro"],
                        "Series": ["SX5E, CAC, DAX, BANKS", "EUR/USD", "Brent, Gold",
                                    "VIX (V2X via SX5E rv)", "DE10Y, FR10Y, IT10Y, ESTR", "EA HICP"],
                        "Coverage": ["2005-2026"] * 6})
    add_table(s, cov, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.8))


def slide_regime_method(prs):
    s = add_slide(prs, "Regime Detection - Methodology", "Phase 3")
    add_bullets(s, [
        "3-state Gaussian Hidden Markov Model fit on 5 standardized macro features",
        "Best-of-5 random seed initialization (defends against local optima)",
        "EONIA/ESTR splice for full-history short-rate continuity (back to 2005)",
        "States ordered by SX5E realized vol -> calm < transition < crisis (interpretable)"
    ], Inches(0.5), Inches(1.3), Inches(12), Inches(2.3), size=14)
    add_text(s, "Feature set:", Inches(0.5), Inches(3.8), Inches(12), Inches(0.4), size=16, bold=True, color=NAVY)
    add_bullets(s, [
        "SX5E_rv_60d - equity stress",
        "SPREAD_IT_DE - Italian sovereign stress",
        "CORR_SX5E_EURUSD_60d - risk-on/risk-off correlation regime",
        "TERM_DE_UNIFIED - ECB stance + growth (uses spliced short rate)",
        "BRENT_rv_60d - commodity stress"
    ], Inches(0.5), Inches(4.2), Inches(12), Inches(2.5), size=13)


def slide_regime_overlay(prs):
    s = add_slide(prs, "Regime Detection - Visual Result", "Phase 3")
    img = REPORTS / "regime_overlay_3state.png"
    if img.exists():
        add_image(s, img, Inches(0.5), Inches(1.4), width=Inches(12.3))
    add_text(s, "Cleanly identifies: 2008-09 GFC; 2011-12 sovereign crisis; Q1 2016 China/oil; "
              "Mar 2020 COVID; Q1 2022 inflation/Russia. The 3-state model surfaces Q1 2016 and the "
              "2007-08 pre-Lehman buildup that a 2-state model misses.",
              Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7), size=12, color=GRAY)


def slide_emission_means(prs):
    s = add_slide(prs, "Regime Detection - Emission Means", "Phase 3")
    em = pd.DataFrame({"calm": [0.164, 1.639, 0.001, 0.177, 0.334],
                       "transition": [0.194, 2.196, 0.036, 1.019, 0.232],
                       "crisis": [0.303, 1.384, 0.181, 1.890, 0.502]},
                      index=["SX5E rv 60d", "IT-DE spread (pp)", "SX5E-EUR/USD corr",
                             "Term spread DE (pp)", "Brent rv 60d"]).reset_index().rename(
                             columns={"index": "Feature"})
    add_table(s, em, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.6))
    add_text(s, "Key insight", Inches(0.5), Inches(4.2), Inches(12), Inches(0.4), size=16, bold=True, color=NAVY)
    add_bullets(s, [
        "Transition has the WIDEST IT-DE spread (220bp) - credit stressed while equity is calm",
        "Crisis has the steepest term spread (189bp) - central bank cutting aggressively into stress",
        "Equity-FX correlation flips from ~0 in calm to +0.18 in crisis - risk-off signature",
        "Brent vol nearly doubles from calm (33%) to crisis (50%) - commodity dislocation"
    ], Inches(0.5), Inches(4.6), Inches(12), Inches(2.5), size=13)


def slide_current_regime(prs):
    s = add_slide(prs, "Current Regime - As of February 2026", "Phase 3")
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.0))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF4, 0xD3, 0x5E); box.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(12), Inches(0.7))
    tb.text_frame.text = "TRANSITION REGIME - 22 weeks running (since Sep 2025)"
    p = tb.text_frame.paragraphs[0]; p.font.size, p.font.bold, p.font.color.rgb = Pt(22), True, NAVY
    add_text(s, "Macro features driving the classification", Inches(0.5), Inches(2.6), Inches(12), Inches(0.4),
              size=16, bold=True, color=NAVY)
    add_bullets(s, [
        "Italian-German 10Y spread elevated relative to post-2022 average",
        "German term spread positive and steepening - ECB has paused, curve normalizing",
        "Equity realized vol contained - ~17% annualized, well below crisis levels",
        "Equity-FX correlation near zero - not yet showing risk-off co-movement"
    ], Inches(0.5), Inches(3.0), Inches(12), Inches(2.5), size=14)
    add_text(s, "Risk implication: macro warning signals are amber while equity markets remain "
              "complacent - exactly the regime where institutional risk teams raise flags.",
              Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.0), size=14, color=NAVY)


def slide_risk_methodology(prs):
    s = add_slide(prs, "Risk Engine - Methodology", "Phase 4")
    add_bullets(s, [
        "Multi-asset portfolio: 50% SX5E + 15% EUR/USD + 20% Brent + 15% Gold",
        "Three VaR variants: historical, parametric, Cornish-Fisher (skew/kurtosis adjusted)",
        "Expected Shortfall: historical and parametric (Gaussian)",
        "Drawdown analysis; rolling beta of BANKS to weekly change in IT-DE spread",
        "Excel parallel build via openpyxl with named ranges - interview-walkable"
    ], Inches(0.5), Inches(1.3), Inches(12), Inches(3.0), size=14)
    add_text(s, "Headline finding", Inches(0.5), Inches(4.7), Inches(12), Inches(0.4), size=16, bold=True, color=NAVY)
    add_text(s, "99% Cornish-Fisher VaR is 11.0% - DOUBLE the parametric Gaussian VaR of 5.0%. "
              "Portfolio skew -0.80 and excess kurtosis 4.4 mean parametric VaR understates 99% "
              "tail risk by more than half. Project-internal evidence of why parametric VaR fails.",
              Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.7), size=14, color=GRAY)


def slide_risk_conditional(prs):
    s = add_slide(prs, "Risk Engine - Regime-Conditional", "Phase 4")
    rc = pd.DataFrame({"calm": [0.124, 0.027, 0.055, 0.042, 0.069, -0.80, 3.38],
                       "transition": [0.128, 0.035, 0.053, 0.046, 0.065, -0.70, 2.13],
                       "crisis": [0.254, 0.049, 0.121, 0.093, 0.175, -1.56, 6.64]},
                      index=["Annualized vol", "Hist VaR 95%", "Hist VaR 99%", "Hist ES 95%",
                             "Hist ES 99%", "Skew", "Excess kurtosis"]).reset_index().rename(
                             columns={"index": "Metric"})
    add_table(s, rc, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.5))
    add_text(s, "The institutional payoff", Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
              size=16, bold=True, color=NAVY)
    add_bullets(s, [
        "Crisis ES 99% is 17.5% - 2.5x the calm ES 99% of 6.9%",
        "Excess kurtosis nearly doubles calm (3.4) to crisis (6.6) - tails get fatter",
        "An allocator using unconditional VaR would mis-price risk by ~2.5x in crisis"
    ], Inches(0.5), Inches(5.4), Inches(12), Inches(2.0), size=14)


def slide_strategies(prs):
    s = add_slide(prs, "Tactical Allocation - Strategies Tested", "Phase 5/5b")
    strat = pd.DataFrame({"Strategy": ["EW", "ERC", "VT-ERC", "Regime-Tilt", "Vol-Forecast", "DD-Predict"],
                          "Description": ["Equal-weight benchmark (25/25/25/25)",
                                           "Equal Risk Contribution (risk parity)",
                                           "ERC scaled to 10% annualized vol target",
                                           "VT-ERC + equity tilt by HMM regime label",
                                           "VT-ERC + equity tilt by LightGBM forecasted vol",
                                           "VT-ERC + gross scaling by LightGBM P(drawdown)"]})
    add_table(s, strat, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.0), size=11)
    add_text(s, "Walk-forward methodology", Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
              size=16, bold=True, color=NAVY)
    add_bullets(s, [
        "LightGBM models refit every 13 weeks on expanding window (no lookahead)",
        "Covariance recomputed at each weekly rebalance",
        "Transaction costs: 5bp on actual one-way turnover",
        "Backtest period: 2010-12 to 2026-05 (~819 weeks)"
    ], Inches(0.5), Inches(5.0), Inches(12), Inches(2.0), size=13)


def slide_strategy_results(prs):
    s = add_slide(prs, "Tactical Allocation - Honest Findings", "Phase 5/5b")
    summ = pd.DataFrame({"EW": [0.0287, 0.1266, 0.287, -0.447, 0.000],
                         "ERC": [0.0192, 0.0899, 0.257, -0.367, 0.156],
                         "VT-ERC": [0.0158, 0.1008, 0.207, -0.402, 0.239],
                         "Regime-Tilt": [0.0127, 0.1054, 0.174, -0.430, 0.305],
                         "Vol-Forecast": [0.0165, 0.1014, 0.213, -0.402, 0.370],
                         "DD-Predict": [0.0108, 0.0994, 0.159, -0.402, 2.057]},
                        index=["Ann return", "Ann vol", "Sharpe", "Max drawdown", "Ann turnover"]
                        ).reset_index().rename(columns={"index": "Metric"})
    add_table(s, summ, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.6))
    add_text(s, "Key finding", Inches(0.5), Inches(4.2), Inches(12), Inches(0.4), size=16, bold=True, color=RED)
    add_bullets(s, [
        "Equal Weight has the HIGHEST Sharpe (0.29) - sophisticated strategies underperformed",
        "Regime-Tilt failed because crisis regimes contain post-crash recovery weeks",
        "DD-Predict failed because weekly drawdown is essentially unpredictable (AUC 0.51)",
        "ERC was the only sophisticated strategy that earned its complexity (-30% vol, similar Sharpe)"
    ], Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), size=13)


def slide_strategy_chart(prs):
    s = add_slide(prs, "Tactical Allocation - Cumulative Returns", "Phase 5/5b")
    img = REPORTS / "strategy_comparison.png"
    if img.exists():
        add_image(s, img, Inches(0.5), Inches(1.4), width=Inches(12.3))
    add_text(s, "EW pulls ahead in the 2024-2026 recovery; sophisticated strategies cluster below. "
              "Most EW outperformance is concentrated in the post-2022 risk-on regime.",
              Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6), size=12, color=GRAY)


def slide_stress(prs):
    s = add_slide(prs, "Stress Testing", "Phase 6")
    add_bullets(s, [
        "8 institutional benchmark scenarios (ECB, oil, FX, equity crash, sovereign stress)",
        "Asset shocks from rolling 156-week empirical betas to each scenario trigger",
        "Direct overrides for portfolio-asset triggers (equity crash hits SX5E directly)",
        "Live Streamlit interface for custom scenario design with combined-shock P&L"
    ], Inches(0.5), Inches(1.3), Inches(12), Inches(2.0), size=14)
    add_text(s, "Sample scenarios", Inches(0.5), Inches(3.5), Inches(12), Inches(0.4), size=16, bold=True, color=NAVY)
    try:
        stress = load_panel("stress_results")
        disp = stress[["scenario", "trigger", "trigger_shock", "portfolio_pnl"]].copy()
        disp.columns = ["Scenario", "Trigger", "Shock", "Portfolio P&L"]
        add_table(s, disp, Inches(0.5), Inches(3.9), Inches(12.3), Inches(3.0), size=10)
    except Exception:
        add_text(s, "Run scripts/build_stress.py to populate stress results.",
                  Inches(0.5), Inches(3.9), Inches(12), Inches(0.5), size=12, color=RED)


def slide_conclusions(prs):
    s = add_slide(prs, "Conclusions", "Wrap-up")
    add_text(s, "What worked", Inches(0.5), Inches(1.3), Inches(12), Inches(0.4), size=18, bold=True, color=NAVY)
    add_bullets(s, [
        "3-state HMM cleanly identifies all major European stress events 2008-2026 including Q1 2016",
        "Regime-conditional risk metrics (ES99: 17.5% crisis vs 6.9% calm) - the institutional payoff",
        "Equal Risk Contribution reduces vol by 30% with similar Sharpe to equal weight",
        "Cornish-Fisher VaR captures fat-tail risk that parametric Gaussian understates by 2x"
    ], Inches(0.5), Inches(1.7), Inches(12), Inches(2.5), size=13)
    add_text(s, "What did not", Inches(0.5), Inches(4.2), Inches(12), Inches(0.4), size=18, bold=True, color=RED)
    add_bullets(s, [
        "Regime-conditional tactical tilts: crisis returns are positive (+20%/yr) due to recoveries",
        "ML drawdown probability forecasting: AUC 0.51 - no exploitable signal at weekly horizon",
        "Conclusion: regime detection is a risk-monitoring tool, not a return-timing signal"
    ], Inches(0.5), Inches(4.6), Inches(12), Inches(2.5), size=13)


def slide_extensions(prs):
    s = add_slide(prs, "Extensions & Limitations", "Wrap-up")
    add_text(s, "Acknowledged limitations", Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
              size=16, bold=True, color=NAVY)
    add_bullets(s, [
        "HMM regime labels for tactical allocation use Phase 3 full-sample fit (mild lookahead)",
        "Stress transmission via empirical betas - assumes linear, time-invariant relationships",
        "Weekly frequency limits ML training sample for tail-event prediction"
    ], Inches(0.5), Inches(1.7), Inches(12), Inches(2.0), size=13)
    add_text(s, "Natural extensions", Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
              size=16, bold=True, color=NAVY)
    add_bullets(s, [
        "Walk-forward HMM refitting (annual expanding window) for fully institutional backtest",
        "Daily-frequency vol surface integration (Eurex SX5E options) for forward-looking vol",
        "Factor-model stress transmission with conditional copulas for tail dependence",
        "Live decision-support: stream regime probabilities into trader dashboards"
    ], Inches(0.5), Inches(4.4), Inches(12), Inches(2.5), size=13)


def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1.5))
    tb.text_frame.text = "Thank you"
    p = tb.text_frame.paragraphs[0]
    p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(60), True, WHITE, PP_ALIGN.CENTER
    sub = s.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(0.6))
    sub.text_frame.text = "Questions?"
    sp = sub.text_frame.paragraphs[0]
    sp.font.size, sp.font.color.rgb, sp.alignment = Pt(20), WHITE, PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_title(prs); slide_summary(prs); slide_data(prs)
    slide_regime_method(prs); slide_regime_overlay(prs); slide_emission_means(prs)
    slide_current_regime(prs); slide_risk_methodology(prs); slide_risk_conditional(prs)
    slide_strategies(prs); slide_strategy_results(prs); slide_strategy_chart(prs)
    slide_stress(prs); slide_conclusions(prs); slide_extensions(prs); slide_thanks(prs)
    out = REPORTS / "macro_regime_methodology_deck.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Deck saved: {out}"); print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
